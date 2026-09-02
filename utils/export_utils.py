# -*- coding: utf-8 -*-
"""
export_utils.py
-----------------
조직도 데이터를 Excel(인원표), PPTX(편집 가능한 조직도 슬라이드),
HTML(진단 보고서)로 내보내는 유틸리티.
"""

from __future__ import annotations
import io
import base64
import datetime as dt

import pandas as pd
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

from utils.org_data import DEPT_TYPE_COLOR, DIAGNOSTIC_CHECKPOINTS
from utils.org_chart_builder import tree_to_dataframe, total_headcount, tree_to_dot

BRAND_NAME = "주식회사 메디엄"


def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ---------------------------------------------------------------------------
# Excel 내보내기
# ---------------------------------------------------------------------------

def build_excel_report(tree: dict, hospital_info: dict, ai_diagnosis: dict | None = None) -> bytes:
    df = tree_to_dataframe(tree)

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        # 1) 조직 인원표
        df.to_excel(writer, sheet_name="조직인원표", index=False)

        # 2) 병원 개요
        overview_rows = [
            {"항목": "병원명", "내용": hospital_info.get("name", "")},
            {"항목": "병원 종별", "내용": hospital_info.get("type", "")},
            {"항목": "병상 수", "내용": hospital_info.get("beds", "")},
            {"항목": "개원 단계", "내용": hospital_info.get("opening_stage", "")},
            {"항목": "총 편성 인원", "내용": total_headcount(tree)},
            {"항목": "작성일", "내용": dt.date.today().isoformat()},
            {"항목": "작성", "내용": BRAND_NAME},
        ]
        pd.DataFrame(overview_rows).to_excel(writer, sheet_name="병원개요", index=False)

        # 3) 진단 체크리스트
        checklist_df = pd.DataFrame({"점검 항목": DIAGNOSTIC_CHECKPOINTS})
        checklist_df.to_excel(writer, sheet_name="조직진단체크리스트", index=False)

        # 4) AI 진단 결과 (있으면)
        if ai_diagnosis:
            issues = ai_diagnosis.get("issues", [])
            recs = ai_diagnosis.get("recommendations", [])
            max_len = max(len(issues), len(recs), 1)
            issues += [""] * (max_len - len(issues))
            recs += [""] * (max_len - len(recs))
            pd.DataFrame({"진단된 문제점": issues, "AI 권장 조치": recs}).to_excel(
                writer, sheet_name="AI진단결과", index=False
            )

        # 서식 적용
        wb = writer.book
        header_fill = PatternFill(start_color="1F2937", end_color="1F2937", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for col_idx, cell in enumerate(ws[1], start=1):
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center", vertical="center")
            for col_cells in ws.columns:
                length = max(len(str(c.value)) if c.value is not None else 0 for c in col_cells)
                col_letter = get_column_letter(col_cells[0].column)
                ws.column_dimensions[col_letter].width = min(max(length * 1.3 + 4, 12), 60)
            ws.freeze_panes = "A2"

    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Graphviz 렌더링 공용 헬퍼
# ---------------------------------------------------------------------------

def render_dot_to_png(dot_str: str, dpi: int = 200) -> bytes | None:
    """DOT 문자열을 PNG 바이트로 렌더링. 서버에 Graphviz 바이너리가 없으면 None."""
    try:
        import graphviz
        src = graphviz.Source(dot_str)
        return src.pipe(format="png")
    except Exception:
        return None


def _png_pixel_size(png_bytes: bytes) -> tuple[int, int]:
    """PNG 헤더에서 (width, height) 픽셀 크기를 추출 (Pillow 의존성 없이)."""
    import struct
    if png_bytes[:8] != b"\x89PNG\r\n\x1a\n":
        return (1600, 900)  # fallback
    width, height = struct.unpack(">II", png_bytes[16:24])
    return (width, height)


# ---------------------------------------------------------------------------
# PPTX 내보내기 — 조직도 슬라이드
# ---------------------------------------------------------------------------
# 조직도 자체는 Graphviz로 렌더링한 고해상도 이미지를 슬라이드에 삽입한다.
# (앱 화면의 조직도 미리보기와 완전히 동일한 배치 · 명단 기입 그리드가 그대로 보존됨)
# 대신 제목·부제·범례·안내 문구는 개별 텍스트 상자로 넣어 PowerPoint에서 자유롭게 수정할 수 있다.

def build_pptx_org_chart(tree: dict, hospital_info: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6))
    tf = title_box.text_frame
    tf.text = f"{hospital_info.get('name', '병원')} 조직도"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.4))
    stf = subtitle_box.text_frame
    stf.text = (f"{hospital_info.get('type', '')} · 병상 {hospital_info.get('beds', '-')}개 · "
                f"{hospital_info.get('opening_stage', '')} · {BRAND_NAME} 작성")
    stf.paragraphs[0].font.size = Pt(12)
    stf.paragraphs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    content_left = Inches(0.4)
    content_top = Inches(1.35)
    content_width = Inches(12.5)
    content_height = Inches(5.15)

    if not tree:
        empty_box = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(0.5))
        empty_box.text_frame.text = "조직도가 비어 있습니다. 앱에서 노드를 추가한 뒤 다시 생성해주세요."
        return _presentation_to_bytes(prs)

    dot_str = tree_to_dot(tree, title="", dpi=200)
    png_bytes = render_dot_to_png(dot_str)

    if png_bytes:
        px_w, px_h = _png_pixel_size(png_bytes)
        aspect = px_w / px_h if px_h else 1.6
        if content_width / aspect <= content_height:
            disp_w = content_width
            disp_h = Emu(int(content_width / aspect))
        else:
            disp_h = content_height
            disp_w = Emu(int(content_height * aspect))
        left = content_left + Emu(int((content_width - disp_w) / 2))
        top = content_top + Emu(int((content_height - disp_h) / 2))
        slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=disp_w, height=disp_h)
    else:
        warn_box = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(0.6))
        warn_box.text_frame.text = (
            "⚠ 조직도 이미지를 생성하지 못했습니다. 서버에 Graphviz 바이너리가 설치되어 있는지 "
            "(packages.txt) 확인해주세요."
        )

    # 범례
    legend_top = Inches(6.75)
    legend_left = Inches(0.4)
    for i, (dept_type, color) in enumerate(DEPT_TYPE_COLOR.items()):
        x = legend_left + Inches(1.72) * i
        box = slide.shapes.add_shape(1, x, legend_top, Inches(0.16), Inches(0.16))
        box.fill.solid()
        box.fill.fore_color.rgb = _hex_to_rgb(color)
        box.line.fill.background()
        box.shadow.inherit = False
        label = slide.shapes.add_textbox(x + Inches(0.22), legend_top - Inches(0.05), Inches(1.5), Inches(0.3))
        label.text_frame.text = dept_type
        label.text_frame.paragraphs[0].font.size = Pt(9)

    note_box = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3))
    note_box.text_frame.text = "※ 실선 상자(직급)의 '성명' 및 명단 그리드(부서)는 인쇄 후 손으로 기입하는 용도입니다."
    note_box.text_frame.paragraphs[0].font.size = Pt(9)
    note_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

    return _presentation_to_bytes(prs)


def _presentation_to_bytes(prs: Presentation) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# HTML 진단 보고서
# ---------------------------------------------------------------------------

def _dot_to_svg_base64(dot_str: str) -> str | None:
    try:
        import graphviz
        src = graphviz.Source(dot_str)
        svg_bytes = src.pipe(format="svg")
        b64 = base64.b64encode(svg_bytes).decode("ascii")
        return f"data:image/svg+xml;base64,{b64}"
    except Exception:
        return None


def build_html_report(tree: dict, hospital_info: dict, ai_diagnosis: dict | None,
                       warnings: list[dict], dot_str: str) -> str:
    df = tree_to_dataframe(tree)
    table_html = df.to_html(index=False, classes="data-table", border=0)

    svg_data_uri = _dot_to_svg_base64(dot_str)
    chart_html = (
        f'<img src="{svg_data_uri}" style="max-width:100%; border-radius:8px;" />'
        if svg_data_uri else
        "<p style='color:#888;'>(조직도 이미지를 생성하려면 서버에 Graphviz 바이너리가 설치되어 있어야 합니다.)</p>"
    )

    diag_html = ""
    if ai_diagnosis:
        issues = "".join(f"<li>{i}</li>" for i in ai_diagnosis.get("issues", []))
        recs = "".join(f"<li>{r}</li>" for r in ai_diagnosis.get("recommendations", []))
        summary = ai_diagnosis.get("summary", "")
        diag_html = f"""
        <div class="card">
          <h2>🤖 AI 조직 진단 요약</h2>
          <p>{summary}</p>
          <div class="two-col">
            <div>
              <h3>진단된 문제점</h3>
              <ul>{issues or '<li>없음</li>'}</ul>
            </div>
            <div>
              <h3>AI 권장 조치</h3>
              <ul>{recs or '<li>없음</li>'}</ul>
            </div>
          </div>
        </div>
        """

    warn_html = ""
    if warnings:
        items = "".join(f"<li>{w['메시지']}</li>" for w in warnings)
        warn_html = f"""
        <div class="card warn">
          <h2>⚠️ 통제범위 경고</h2>
          <ul>{items}</ul>
        </div>
        """

    checklist_html = "".join(f"<li>{c}</li>" for c in DIAGNOSTIC_CHECKPOINTS)

    today = dt.date.today().isoformat()

    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8" />
<title>{hospital_info.get('name', '병원')} 조직도 진단 보고서</title>
<style>
  body {{ font-family: 'Malgun Gothic', 'Apple SD Gothic Neo', sans-serif; background:#0F172A; color:#E2E8F0; margin:0; padding:0; }}
  .wrap {{ max-width: 980px; margin: 0 auto; padding: 40px 24px 80px; }}
  header {{ border-bottom: 1px solid #1E293B; padding-bottom: 24px; margin-bottom: 32px; }}
  header .brand {{ color:#D4AF37; font-weight:700; letter-spacing:1px; font-size:13px; }}
  header h1 {{ font-size: 30px; margin: 8px 0 4px; color:#FFFFFF; }}
  header .meta {{ color:#94A3B8; font-size: 14px; }}
  .card {{ background:#111827; border:1px solid #1E293B; border-radius:14px; padding:28px; margin-bottom:24px; }}
  .card.warn {{ border-color:#7C2D12; background:#1C1410; }}
  h2 {{ font-size:19px; color:#F1F5F9; margin-top:0; }}
  h3 {{ font-size:15px; color:#CBD5E1; }}
  .two-col {{ display:flex; gap:32px; flex-wrap:wrap; }}
  .two-col > div {{ flex:1; min-width:260px; }}
  ul {{ padding-left: 20px; line-height:1.7; color:#CBD5E1; }}
  .data-table {{ width:100%; border-collapse: collapse; font-size: 13px; }}
  .data-table th {{ background:#1F2937; color:#F8FAFC; text-align:left; padding:8px 10px; }}
  .data-table td {{ padding:7px 10px; border-bottom:1px solid #1E293B; color:#E2E8F0; }}
  .chart-box {{ text-align:center; background:#0B1220; border-radius:10px; padding:16px; }}
  footer {{ text-align:center; color:#475569; font-size:12px; margin-top:40px; }}
</style>
</head>
<body>
<div class="wrap">
  <header>
    <div class="brand">{BRAND_NAME} · HOSPITAL ORGANIZATION DESIGN REPORT</div>
    <h1>{hospital_info.get('name', '병원')} 조직도 진단 보고서</h1>
    <div class="meta">{hospital_info.get('type','')} · 병상 {hospital_info.get('beds','-')}개 ·
      {hospital_info.get('opening_stage','')} · 작성일 {today} · 총 편성인원 {total_headcount(tree)}명</div>
  </header>

  <div class="card">
    <h2>🌳 조직도</h2>
    <div class="chart-box">{chart_html}</div>
  </div>

  {diag_html}
  {warn_html}

  <div class="card">
    <h2>📋 조직 인원표</h2>
    {table_html}
  </div>

  <div class="card">
    <h2>✅ 조직 진단 체크리스트 (일반 참고용)</h2>
    <ul>{checklist_html}</ul>
    <p style="color:#64748B; font-size:12px; margin-top:16px;">
      ※ 본 보고서의 조직 구조 및 인력 배치는 일반적인 컨설팅 참고 프레임워크이며,
      의료법 시행규칙상 정확한 법정 정원·필수 인력 기준은 최신 법령 및 관할 보건소 확인이 필요합니다.
    </p>
  </div>

  <footer>Generated by {BRAND_NAME} · 병원 조직도 설계 AI</footer>
</div>
</body>
</html>"""
